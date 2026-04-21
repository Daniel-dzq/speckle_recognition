#!/usr/bin/env python3
"""
Export the 26 letter slides from the PowerPoint file as PNG images.

Usage:
    python scripts/export_letter_images.py

This script tries two methods in order:
  1. PowerPoint COM automation (Windows, requires MS Office installed)
  2. Pillow-based software rendering (cross-platform fallback)

Output images are saved to:
    <project_root>/letter_images/A.png  ...  Z.png

The SLMWindow will automatically use these images instead of rendering
text when they exist.
"""

import os
import sys

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPTX_PATH = os.path.join(ROOT, "26个英文字母（大）.pptx")
OUT_DIR = os.path.join(ROOT, "letter_images")

TARGET_W = 1920
TARGET_H = 1080


def export_via_com(pptx_path: str, out_dir: str) -> bool:
    """Use PowerPoint COM automation (Windows only) to export slides as PNG."""
    try:
        import comtypes.client  # type: ignore
    except ImportError:
        print("[COM] comtypes not installed, skipping COM export.")
        return False

    if sys.platform != "win32":
        return False

    print("[COM] Attempting PowerPoint COM export...")
    pptx_abs = os.path.abspath(pptx_path)
    out_abs = os.path.abspath(out_dir)
    os.makedirs(out_abs, exist_ok=True)

    ppt = None
    prs = None
    try:
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        ppt.Visible = 1
        prs = ppt.Presentations.Open(pptx_abs, ReadOnly=True, Untitled=False, WithWindow=False)

        slides_list = list(prs.Slides)
        from pptx import Presentation as PyPptx
        meta = PyPptx(pptx_path)
        meta_slides = list(meta.slides)

        for i, slide in enumerate(slides_list):
            letter = "A"
            if i < len(meta_slides):
                for shape in meta_slides[i].shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip().upper()
                        if t and t.isalpha():
                            letter = t[0]
                            break

            out_path = os.path.join(out_abs, f"{letter}.png")
            slide.Export(out_path, "PNG", TARGET_W, TARGET_H)
            print(f"  Slide {i+1} -> {letter}.png")

        print(f"[COM] Exported {len(slides_list)} slides to {out_abs}")
        return True

    except Exception as e:
        print(f"[COM] Export failed: {e}")
        return False

    finally:
        try:
            if prs is not None:
                prs.Close()
        except Exception:
            pass
        try:
            if ppt is not None:
                ppt.Quit()
        except Exception:
            pass


def export_via_pillow(pptx_path: str, out_dir: str) -> bool:
    """Render slides with Pillow using the same styling as the PPTX."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("[Pillow] Pillow not installed. Run: pip install Pillow")
        return False

    try:
        from pptx import Presentation
    except ImportError:
        print("[Pillow] python-pptx not installed. Run: pip install python-pptx")
        return False

    print("[Pillow] Rendering letter images with Pillow...")
    os.makedirs(out_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    slides_list = list(prs.slides)

    # Slide dimensions in EMU; convert to pixel ratio
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    # Resolve scheme colors: tx1=black (bg), bg1=white (text)
    bg_color = (0, 0, 0)
    fg_color = (255, 255, 255)

    # Font: Calibri Bold from theme. Try to find a system font.
    font_candidates = [
        "calibrib.ttf",   # Calibri Bold (Windows)
        "Calibri Bold.ttf",
        "calibri.ttf",
        "arialbd.ttf",    # Arial Bold fallback
        "arial.ttf",
    ]

    font_dir_candidates = [
        r"C:\Windows\Fonts",
        "/usr/share/fonts/truetype",
        "/System/Library/Fonts",
        os.path.expanduser("~/Library/Fonts"),
    ]

    font_path = None
    for fd in font_dir_candidates:
        if not os.path.isdir(fd):
            continue
        for fc in font_candidates:
            fp = os.path.join(fd, fc)
            if os.path.isfile(fp):
                font_path = fp
                break
        if font_path:
            break

    if font_path:
        print(f"[Pillow] Using font: {font_path}")
    else:
        print("[Pillow] No TTF font found; using Pillow default (will look different).")

    for i, slide in enumerate(slides_list):
        letter = chr(ord("A") + i)
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip().upper()
                if t and t.isalpha():
                    letter = t[0]
                    break

        img = Image.new("RGB", (TARGET_W, TARGET_H), bg_color)
        draw = ImageDraw.Draw(img)

        # Binary-search for the largest font size that fits inside TARGET_H * 0.9
        max_h = int(TARGET_H * 0.90)
        max_w = int(TARGET_W * 0.90)
        lo, hi = 10, TARGET_H
        best_font = None
        best_size = lo
        while lo <= hi:
            mid = (lo + hi) // 2
            if font_path:
                try:
                    f = ImageFont.truetype(font_path, mid)
                except Exception:
                    f = ImageFont.load_default()
            else:
                f = ImageFont.load_default()

            bbox = draw.textbbox((0, 0), letter, font=f)
            tw = bbox[2] - bbox[0]
            th = bbox[3] - bbox[1]
            if tw <= max_w and th <= max_h:
                best_font = f
                best_size = mid
                lo = mid + 1
            else:
                hi = mid - 1

        if best_font is None:
            if font_path:
                best_font = ImageFont.truetype(font_path, 10)
            else:
                best_font = ImageFont.load_default()

        bbox = draw.textbbox((0, 0), letter, font=best_font)
        tw = bbox[2] - bbox[0]
        th = bbox[3] - bbox[1]
        x = (TARGET_W - tw) // 2 - bbox[0]
        y = (TARGET_H - th) // 2 - bbox[1]
        draw.text((x, y), letter, fill=fg_color, font=best_font)

        out_path = os.path.join(out_dir, f"{letter}.png")
        img.save(out_path, "PNG")
        print(f"  Slide {i+1} -> {letter}.png  (font size {best_size}px)")

    print(f"[Pillow] Exported {len(slides_list)} images to {out_dir}")
    return True


def main():
    if not os.path.isfile(PPTX_PATH):
        print(f"[ERROR] PPTX not found: {PPTX_PATH}")
        sys.exit(1)

    os.makedirs(OUT_DIR, exist_ok=True)

    # Try COM first (highest fidelity on Windows), then fall back to Pillow
    if not export_via_com(PPTX_PATH, OUT_DIR):
        if not export_via_pillow(PPTX_PATH, OUT_DIR):
            print("[ERROR] Both export methods failed.")
            sys.exit(1)

    print(f"\nDone! Letter images are in: {OUT_DIR}")
    print("Restart the demo GUI – the SLM window will now use these images.")


if __name__ == "__main__":
    main()
