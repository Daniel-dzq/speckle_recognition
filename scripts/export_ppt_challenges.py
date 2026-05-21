#!/usr/bin/env python3
"""
Export challenge slides from a PowerPoint deck to PNG images and manifest.json.

Usage:
    python scripts/export_ppt_challenges.py --input input.pptx --output_dir challenge_inputs

Methods (in order):
  1. LibreOffice headless (highest fidelity, if installed)
  2. python-pptx + Pillow (text render + embedded picture extract)
"""

from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from io import BytesIO
from typing import List, Optional, Tuple

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)

TARGET_W = 1920
TARGET_H = 1080
DEFAULT_AVATAR_LABELS = ("boy", "girl")


def parse_args():
    p = argparse.ArgumentParser(description="Export PPT challenge slides to PNG + manifest")
    p.add_argument("--input", type=str, default=os.path.join(ROOT, "input.pptx"),
                   help="Source .pptx file")
    p.add_argument("--output_dir", type=str, default=os.path.join(ROOT, "challenge_inputs"),
                   help="Output directory for PNGs and manifest.json")
    p.add_argument("--width", type=int, default=TARGET_W)
    p.add_argument("--height", type=int, default=TARGET_H)
    p.add_argument("--avatar_labels", type=str, default="boy,girl",
                   help="Comma-separated labels for picture-only slides (in order)")
    p.add_argument("--force", action="store_true",
                   help="Overwrite existing PNG files")
    return p.parse_args()


def _sanitize_filename(label: str, slide_index: int) -> str:
    """Filesystem-safe stem; fall back to slide_XX if label is empty or unsafe."""
    raw = (label or "").strip()
    if raw:
        safe = re.sub(r"[^\w.\-]+", "_", raw, flags=re.UNICODE).strip("._")
        if safe:
            return safe
    return f"slide_{slide_index:02d}"


def _find_libreoffice() -> Optional[str]:
    candidates = [
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for path in candidates:
        if path and os.path.isfile(path) and os.access(path, os.X_OK):
            return path
    return None


def _extract_slide_label(slide, slide_index: int, avatar_queue: List[str]) -> Tuple[str, str]:
    """
    Return (label, method) where method is 'text' or 'avatar_fallback' or 'slide_index'.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    texts: List[str] = []
    has_picture = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_picture = True

    if texts:
        primary = texts[0].split("\n")[0].strip()
        if primary:
            return primary, "text"

    if has_picture and avatar_queue:
        return avatar_queue.pop(0), "avatar_fallback"

    return f"slide_{slide_index:02d}", "slide_index"


def export_via_libreoffice(
    pptx_path: str,
    out_dir: str,
    width: int,
    height: int,
    avatar_labels: List[str],
) -> bool:
    """Export slides with LibreOffice; requires separate labeling pass via pptx metadata."""
    soffice = _find_libreoffice()
    if not soffice:
        return False

    pptx_abs = os.path.abspath(pptx_path)
    out_abs = os.path.abspath(out_dir)
    os.makedirs(out_abs, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="ppt_export_") as tmp:
        cmd = [
            soffice,
            "--headless",
            "--norestore",
            "--convert-to", "png",
            "--outdir", tmp,
            pptx_abs,
        ]
        print(f"[LibreOffice] Running: {' '.join(cmd)}")
        try:
            subprocess.run(cmd, check=True, capture_output=True, text=True, timeout=180)
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError) as exc:
            print(f"[LibreOffice] Export failed: {exc}")
            return False

        pngs = sorted(
            f for f in os.listdir(tmp)
            if f.lower().endswith(".png")
        )
        if not pngs:
            print("[LibreOffice] No PNG files produced.")
            return False

        try:
            from pptx import Presentation
        except ImportError:
            print("[LibreOffice] python-pptx required to label slides after LO export.")
            return False

        from PIL import Image

        prs = Presentation(pptx_abs)
        slides = list(prs.slides)
        avatar_queue = list(avatar_labels)
        challenges = []

        for i, png_name in enumerate(pngs):
            slide_index = i + 1
            slide = slides[i] if i < len(slides) else None
            if slide is not None:
                label, _ = _extract_slide_label(slide, slide_index, avatar_queue)
            else:
                label = f"slide_{slide_index:02d}"

            stem = _sanitize_filename(label, slide_index)
            src = os.path.join(tmp, png_name)
            dst = os.path.join(out_abs, f"{stem}.png")
            _resize_png_to_canvas(src, dst, width, height)
            rel = os.path.relpath(dst, ROOT).replace("\\", "/")
            challenges.append({
                "label": label,
                "image": rel,
                "source_slide": slide_index,
                "export_method": "libreoffice",
            })
            print(f"  Slide {slide_index} -> {stem}.png  (label={label!r})")

        _write_manifest(out_abs, pptx_abs, challenges, "libreoffice")
        print(f"[LibreOffice] Exported {len(challenges)} slides to {out_abs}")
        return True


def _resize_png_to_canvas(src_path: str, dst_path: str, width: int, height: int) -> None:
    from PIL import Image

    img = Image.open(src_path)
    if img.mode not in ("RGB", "RGBA"):
        img = img.convert("RGB")

    bg = (0, 0, 0)
    if img.mode == "RGBA":
        canvas = Image.new("RGB", (width, height), bg)
        fitted = _fit_image(img, width, height)
        ox = (width - fitted.width) // 2
        oy = (height - fitted.height) // 2
        canvas.paste(fitted, (ox, oy), fitted)
        canvas.save(dst_path, "PNG")
    else:
        canvas = Image.new("RGB", (width, height), bg)
        fitted = _fit_image(img.convert("RGB"), width, height)
        ox = (width - fitted.width) // 2
        oy = (height - fitted.height) // 2
        canvas.paste(fitted, (ox, oy))
        canvas.save(dst_path, "PNG")


def _fit_image(img, max_w: int, max_h: int):
    from PIL import Image

    w, h = img.size
    scale = min(max_w / w, max_h / h, 1.0)
    nw = max(1, int(w * scale))
    nh = max(1, int(h * scale))
    return img.resize((nw, nh), Image.Resampling.LANCZOS)


def _scheme_colors() -> Tuple[Tuple[int, int, int], Tuple[int, int, int]]:
    return (0, 0, 0), (255, 255, 255)


def _resolve_font_path() -> Optional[str]:
    candidates = [
        "calibrib.ttf",
        "Calibri Bold.ttf",
        "calibri.ttf",
        "arialbd.ttf",
        "Arial Bold.ttf",
        "Helvetica.ttc",
    ]
    dirs = [
        r"C:\Windows\Fonts",
        "/System/Library/Fonts/Supplemental",
        "/System/Library/Fonts",
        "/Library/Fonts",
        os.path.expanduser("~/Library/Fonts"),
        "/usr/share/fonts/truetype",
    ]
    for d in dirs:
        if not os.path.isdir(d):
            continue
        for name in candidates:
            path = os.path.join(d, name)
            if os.path.isfile(path):
                return path
    return None


def _render_text_slide(text: str, out_path: str, width: int, height: int, font_path: Optional[str]) -> None:
    from PIL import Image, ImageDraw, ImageFont

    bg_color, fg_color = _scheme_colors()
    display = text.strip()
    if not display:
        display = "?"

    img = Image.new("RGB", (width, height), bg_color)
    draw = ImageDraw.Draw(img)
    max_h = int(height * 0.90)
    max_w = int(width * 0.90)
    lo, hi = 10, height
    best_font = None
    while lo <= hi:
        mid = (lo + hi) // 2
        try:
            f = ImageFont.truetype(font_path, mid) if font_path else ImageFont.load_default()
        except OSError:
            f = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), display, font=f)
        tw = bbox[2] - bbox[0]
        th = bbox[3] - bbox[1]
        if tw <= max_w and th <= max_h:
            best_font = f
            lo = mid + 1
        else:
            hi = mid - 1

    if best_font is None:
        best_font = ImageFont.truetype(font_path, 48) if font_path else ImageFont.load_default()

    bbox = draw.textbbox((0, 0), display, font=best_font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    x = (width - tw) // 2 - bbox[0]
    y = (height - th) // 2 - bbox[1]
    draw.text((x, y), display, fill=fg_color, font=best_font)
    img.save(out_path, "PNG")


def _export_picture_slide(shape, out_path: str, width: int, height: int) -> None:
    from PIL import Image

    blob = shape.image.blob
    img = Image.open(BytesIO(blob))
    if img.mode not in ("RGB", "RGBA"):
        img = img.convert("RGBA" if "A" in img.mode else "RGB")

    bg = (0, 0, 0)
    canvas = Image.new("RGB", (width, height), bg)
    if img.mode == "RGBA":
        fitted = _fit_image(img, width, height)
        ox = (width - fitted.width) // 2
        oy = (height - fitted.height) // 2
        canvas.paste(fitted, (ox, oy), fitted)
    else:
        fitted = _fit_image(img.convert("RGB"), width, height)
        ox = (width - fitted.width) // 2
        oy = (height - fitted.height) // 2
        canvas.paste(fitted, (ox, oy))
    canvas.save(out_path, "PNG")


def export_via_pptx(
    pptx_path: str,
    out_dir: str,
    width: int,
    height: int,
    avatar_labels: List[str],
    force: bool,
) -> bool:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("[ERROR] python-pptx is not installed. Run: pip install python-pptx")
        return False

    try:
        from PIL import Image  # noqa: F401
    except ImportError:
        print("[ERROR] Pillow is not installed. Run: pip install Pillow")
        return False

    print("[pptx] Exporting slides with python-pptx + Pillow...")
    os.makedirs(out_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    font_path = _resolve_font_path()
    if font_path:
        print(f"[pptx] Font: {font_path}")
    else:
        print("[pptx] No system TTF found; using default font for text slides.")

    avatar_queue = list(avatar_labels)
    challenges = []
    used_stems: dict[str, int] = {}

    for i, slide in enumerate(slides, start=1):
        label, method = _extract_slide_label(slide, i, avatar_queue)
        stem = _sanitize_filename(label, i)
        if stem in used_stems:
            used_stems[stem] += 1
            stem = f"{stem}_{used_stems[stem]}"
        else:
            used_stems[stem] = 1

        out_path = os.path.join(out_dir, f"{stem}.png")
        if os.path.isfile(out_path) and not force:
            print(f"  Slide {i}: skip existing {stem}.png")
        else:
            picture_shape = None
            text_value = None
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and picture_shape is None:
                    picture_shape = shape
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t and text_value is None:
                        text_value = t.split("\n")[0].strip()

            if picture_shape is not None:
                _export_picture_slide(picture_shape, out_path, width, height)
                print(f"  Slide {i} -> {stem}.png  (picture, label={label!r})")
            elif text_value:
                _render_text_slide(text_value, out_path, width, height, font_path)
                print(f"  Slide {i} -> {stem}.png  (text={text_value!r})")
            else:
                _render_text_slide(label, out_path, width, height, font_path)
                print(f"  Slide {i} -> {stem}.png  (fallback label={label!r})")

        rel = os.path.relpath(out_path, ROOT).replace("\\", "/")
        challenges.append({
            "label": label,
            "image": rel,
            "source_slide": i,
            "label_method": method,
            "export_method": "pptx",
        })

    _write_manifest(out_dir, pptx_path, challenges, "pptx")
    print(f"[pptx] Exported {len(challenges)} slides to {out_dir}")
    return True


def _write_manifest(out_dir: str, pptx_path: str, challenges: list, method: str) -> None:
    manifest = {
        "source_pptx": os.path.abspath(pptx_path),
        "export_method": method,
        "width": TARGET_W,
        "height": TARGET_H,
        "challenges": challenges,
    }
    path = os.path.join(out_dir, "manifest.json")
    with open(path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2)
    print(f"[manifest] Wrote {path}")


def main():
    args = parse_args()
    pptx_path = os.path.abspath(args.input)
    out_dir = os.path.abspath(args.output_dir)

    if not os.path.isfile(pptx_path):
        print(f"[ERROR] Input file not found: {pptx_path}")
        sys.exit(1)

    avatar_labels = [x.strip() for x in args.avatar_labels.split(",") if x.strip()]

    if export_via_libreoffice(pptx_path, out_dir, args.width, args.height, avatar_labels):
        return

    lo = _find_libreoffice()
    if lo is None:
        print(
            "[INFO] LibreOffice not found. Install LibreOffice for pixel-perfect export, or "
            "use the built-in python-pptx exporter below."
        )
        print(
            "  macOS example: brew install --cask libreoffice\n"
            "  Then run: /Applications/LibreOffice.app/Contents/MacOS/soffice --headless "
            "--convert-to png --outdir <tmpdir> input.pptx"
        )
    else:
        print("[INFO] LibreOffice export failed; falling back to python-pptx.")

    if not export_via_pptx(
        pptx_path, out_dir, args.width, args.height, avatar_labels, args.force,
    ):
        print("[ERROR] Export failed. Install dependencies:")
        print("  pip install python-pptx Pillow")
        sys.exit(1)

    print(f"\nDone. Challenge set ready in: {out_dir}")
    print("Start the demo GUI; it will auto-load challenge_inputs/manifest.json if present.")


if __name__ == "__main__":
    main()
