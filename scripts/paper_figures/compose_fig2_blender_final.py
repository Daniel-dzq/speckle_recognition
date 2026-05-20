#!/usr/bin/env python3
"""
Compose Fig.2 final figure from Blender renders + vector overlays.

Inputs (from Blender pipeline):
  Fig2_optical_setup_9cm_render.png
  Fig2_optical_setup_9cm_polish_closeup.png

Outputs:
  Fig2_optical_setup_9cm.png / .pdf / .svg  — full composed figure
  Fig2_optical_setup_9cm_semieditable.pptx — Blender renders as pictures + editable text/shapes
  Fig2_optical_setup_9cm_annotations.svg   — optional text/line overlay (same layout, smaller)

Usage:
  python3 scripts/paper_figures/compose_fig2_blender_final.py
"""
from __future__ import annotations

import sys
from io import BytesIO
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from matplotlib.colorbar import ColorbarBase
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from PIL import Image

SCRIPT_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRIPT_DIR.parents[1]
OUT = REPO_ROOT / "figures" / "paper" / "Fig2_setup"
BASE = "Fig2_optical_setup_9cm"

C_LABEL = "#1A2744"
C_RED = "#C62828"
C_GREEN = "#2E7D32"
C_CORE = "#90CAF9"
C_CLAD = "#CFD8DC"
C_PURPLE = "#7B1FA2"
C_GREEN_BOX = "#388E3C"


def load_rgb(path: Path) -> np.ndarray:
    im = Image.open(path).convert("RGB")
    return np.asarray(im) / 255.0


def draw_panel_c(ax, rng: np.random.Generator):
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.axis("off")
    box = FancyBboxPatch(
        (4, 4), 92, 92, boxstyle="round,pad=0,rounding_size=2",
        facecolor="none", edgecolor=C_GREEN_BOX, linewidth=1.8, linestyle=(0, (5, 4)),
    )
    ax.add_patch(box)
    ax.text(50, 92, "SLM Letter Image Input (Examples)", ha="center", fontsize=9, fontweight="bold", color=C_LABEL)
    letters = [("A", 10), ("B", 30), ("C", 50), ("Z", 70)]
    for ch, x0 in letters:
        ax.add_patch(mpatches.Rectangle((x0, 55), 17, 20, facecolor="black", edgecolor="#333"))
        ax.text(x0 + 8.5, 65, ch, ha="center", va="center", fontsize=14, color="white", fontweight="bold")
        yy, xx = np.mgrid[0:16, 0:16]
        ph = np.sin(xx * 0.45 + yy * 0.38 + rng.normal(0, 0.15, (16, 16))) * 0.5
        ax.imshow(ph, extent=(x0, x0 + 17, 36, 52), cmap="viridis", origin="lower", vmin=-1, vmax=1)
    ax.text(50, 8, "26 challenge patterns (A–Z)", ha="center", fontsize=8, color=C_LABEL)
    ax.text(5, 95, "(c)", fontsize=11, fontweight="bold", color=C_LABEL, va="top")


def add_scale_bar(fig, ax_main, y_frac_below=0.06):
    """Brackets under main image in axes fraction coordinates."""
    xa, ya, wa, ha = ax_main.get_position().bounds
    y0 = ya - 0.055
    x0, x1 = xa + 0.12, xa + wa - 0.08
    # Proportions 5:1:3
    t5, t1, t3 = 5 / 9, 1 / 9, 3 / 9
    p0 = x0
    p1 = x0 + t5 * (x1 - x0)
    p2 = p1 + t1 * (x1 - x0)
    p3 = x1
    for (a, b, lab) in ((p0, p1, "5 cm"), (p1, p2, "1 cm"), (p2, p3, "3 cm")):
        fig.add_artist(plt.Line2D([a, b], [y0, y0], color=C_LABEL, lw=1.2, transform=fig.transFigure, clip_on=False))
        fig.add_artist(plt.Line2D([a, a], [y0 - 0.008, y0 + 0.008], color=C_LABEL, lw=1, transform=fig.transFigure, clip_on=False))
        fig.add_artist(plt.Line2D([b, b], [y0 - 0.008, y0 + 0.008], color=C_LABEL, lw=1, transform=fig.transFigure, clip_on=False))
        fig.text((a + b) / 2, y0 - 0.018, lab, ha="center", va="top", fontsize=7.5, color=C_LABEL, transform=fig.transFigure)
    fig.text((p0 + p3) / 2, y0 - 0.038, "Total: 9 cm", ha="center", va="top", fontsize=8, fontweight="bold", color=C_LABEL, transform=fig.transFigure)


def legend_box(fig, rect):
    axl = fig.add_axes(rect)
    axl.set_xlim(0, 1)
    axl.set_ylim(0, 1)
    axl.axis("off")
    axl.add_patch(FancyBboxPatch((0.04, 0.06), 0.92, 0.88, boxstyle="round,pad=0.02", facecolor="#F5F5F5", edgecolor="#B0BEC5", lw=1))
    axl.plot([0.1, 0.2], [0.82, 0.82], color=C_RED, lw=3.5, solid_capstyle="round")
    axl.text(0.24, 0.82, "Red channel", fontsize=8, fontweight="bold", color=C_LABEL, va="center")
    axl.text(0.24, 0.72, "(Reference, end-face axial)", fontsize=7, color=C_LABEL, va="center")
    axl.plot([0.1, 0.2], [0.56, 0.56], color="#4CAF50", lw=3.5, solid_capstyle="round")
    axl.text(0.24, 0.56, "Green channel", fontsize=8, fontweight="bold", color=C_LABEL, va="center")
    axl.text(0.24, 0.46, "(Signal, side illumination)", fontsize=7, color=C_LABEL, va="center")
    axl.add_patch(mpatches.Rectangle((0.08, 0.28), 0.12, 0.08, facecolor=C_CORE, edgecolor="#1565C0", lw=0.5))
    axl.text(0.24, 0.32, "POF Core (Higher n)", fontsize=7.5, color=C_LABEL, va="center")
    axl.add_patch(mpatches.Rectangle((0.08, 0.12), 0.12, 0.08, facecolor=C_CLAD, edgecolor="#90A4AE", lw=0.5))
    axl.text(0.24, 0.16, "POF Cladding (Lower n)", fontsize=7.5, color=C_LABEL, va="center")


def compose_figure():
    main_path = OUT / f"{BASE}_render.png"
    pol_path = OUT / f"{BASE}_polish_closeup.png"
    if not main_path.is_file():
        raise FileNotFoundError(main_path)
    if not pol_path.is_file():
        raise FileNotFoundError(pol_path)

    main_img = load_rgb(main_path)
    pol_img = load_rgb(pol_path)
    rng = np.random.default_rng(7)

    fig = plt.figure(figsize=(11.5, 8.4), facecolor="white", dpi=150)
    fig.patch.set_facecolor("white")

    ax_m = fig.add_axes([0.07, 0.30, 0.62, 0.58])
    ax_m.imshow(main_img, aspect="equal")
    ax_m.axis("off")
    ax_m.text(0.02, 0.98, "(a)", transform=ax_m.transAxes, fontsize=11, fontweight="bold", color=C_LABEL, va="top")

    # Labels on main (approximate positions — tune after viewing render)
    def mt(x, y, s, fs=7, ha="left", bold=False):
        ax_m.text(x, y, s, transform=ax_m.transAxes, fontsize=fs, color=C_LABEL, ha=ha, fontweight="bold" if bold else "normal")

    mt(0.01, 0.08, "Red Laser (650 nm)", 8, bold=True)
    mt(0.01, 0.03, "Red channel:\nhorizontal end-face\nreference input", 6.5)
    mt(0.44, 0.92, "Step-index POF\n(9 cm rigid rod)", 7.5, ha="center")
    mt(0.02, 0.14, "Left End-Face\n(Axial Reference Input)", 6.5)
    mt(0.78, 0.14, "Right End-Face\n(Output)", 6.5, ha="center")
    mt(0.52, 0.88, "Side-polished Region\n(Side Illumination / Challenge Input)", 6.5, ha="center")
    mt(0.12, 0.95, "Green Laser (520 nm)", 8, bold=True)
    mt(0.12, 0.90, "Green channel:\noblique side-illumination\nchallenge input", 6.5)
    mt(0.28, 0.72, "Beam Expander", 6.5, ha="center")
    mt(0.36, 0.62, "Collimating Lens", 6.5, ha="center")
    mt(0.48, 0.52, "SLM\n(Spatial Light Modulator)", 6.5, ha="center")
    mt(0.62, 0.44, "Focusing Lens", 6.5, ha="center")
    mt(0.82, 0.62, "CMOS Camera", 7.5, ha="center", bold=True)
    mt(0.84, 0.38, "Speckle\nOutput", 6.5, ha="center")

    add_scale_bar(fig, ax_m)

    ax_b = fig.add_axes([0.06, 0.07, 0.28, 0.20])
    ax_b.imshow(pol_img, aspect="equal")
    ax_b.axis("off")
    ax_b.add_patch(
        FancyBboxPatch(
            (-0.02, -0.02), 1.04, 1.04, transform=ax_b.transAxes,
            facecolor="none", edgecolor=C_PURPLE, linewidth=2, linestyle=(0, (6, 3)),
            clip_on=False,
        )
    )
    ax_b.set_title("Magnified View: Side-polished Region", fontsize=8.5, fontweight="bold", color=C_LABEL, pad=3)
    ax_b.text(0.5, 0.12, "Cladding\n(Lower n)", transform=ax_b.transAxes, ha="center", fontsize=6.5, color=C_LABEL)
    ax_b.text(0.78, 0.55, "Core\n(Higher n)", transform=ax_b.transAxes, ha="center", fontsize=6.5, color=C_LABEL)
    ax_b.text(0.52, 0.88, "Polishing window\n(enhanced side coupling)", transform=ax_b.transAxes, ha="center", fontsize=6.5, color=C_LABEL)
    ax_b.text(0.62, 0.5, "Side-illuminated\nscattering volume", transform=ax_b.transAxes, fontsize=6.5, color=C_LABEL)
    ax_b.text(0.02, 0.98, "(b)", transform=ax_b.transAxes, fontsize=11, fontweight="bold", va="top", color=C_LABEL)

    ax_c = fig.add_axes([0.70, 0.55, 0.27, 0.30])
    draw_panel_c(ax_c, rng)
    cax = fig.add_axes([0.955, 0.58, 0.012, 0.22])
    ColorbarBase(cax, cmap=plt.cm.viridis, orientation="vertical")
    cax.set_title("Phase", fontsize=7, color=C_LABEL, pad=4)
    cax.set_ylabel("2π", fontsize=6.5, color=C_LABEL)

    fa = FancyArrowPatch(
        (0.72, 0.72), (0.48, 0.58), transform=fig.transFigure,
        arrowstyle="-", linestyle=(0, (4, 3)), linewidth=1.2, color="#4CAF50", mutation_scale=12,
    )
    fig.add_artist(fa)

    legend_box(fig, [0.68, 0.07, 0.28, 0.18])

    OUT.mkdir(parents=True, exist_ok=True)
    for ext in ("png", "pdf", "svg"):
        fig.savefig(OUT / f"{BASE}.{ext}", format=ext, dpi=300 if ext == "png" else None, facecolor="white", bbox_inches="tight", pad_inches=0.06)
    plt.close(fig)
    print("Wrote", OUT / f"{BASE}.png", "+ pdf/svg")


def _compose_annotations_svg_only():
    """Lightweight vector-only overlay (scale + legend) for manual alignment in Inkscape."""
    fig = plt.figure(figsize=(11.5, 8.4), facecolor="none")
    fig.patch.set_alpha(0)
    legend_box(fig, [0.68, 0.07, 0.28, 0.18])
    fig.savefig(OUT / f"{BASE}_annotations.svg", facecolor="none", bbox_inches="tight", pad_inches=0.06)
    plt.close(fig)


def build_pptx_semieditable():
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Inches, Pt
    except ImportError:
        print("pip install python-pptx", file=sys.stderr)
        return

    def emu(i: float) -> int:
        return int(round(i * 914400))

    prs = Presentation()
    prs.slide_width = emu(13.333)
    prs.slide_height = emu(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    main_path = OUT / f"{BASE}_render.png"
    pol_path = OUT / f"{BASE}_polish_closeup.png"
    # Background white
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    # Pictures (locked optics)
    slide.shapes.add_picture(str(main_path), emu(0.45), emu(1.9), width=emu(8.35))
    slide.shapes.add_picture(str(pol_path), emu(0.45), emu(0.52), width=emu(3.65))

    def tb(l, t, w, h, txt, pt=8.5, bold=False, align=PP_ALIGN.LEFT):
        bx = slide.shapes.add_textbox(emu(l), emu(t), emu(w), emu(h))
        p = bx.text_frame.paragraphs[0]
        p.text = txt
        p.font.size = Pt(pt)
        p.font.name = "Arial"
        p.font.bold = bold
        p.font.color.rgb = RGBColor(26, 39, 68)
        p.alignment = align
        return bx

    # Panel (a) labels — approximate slide inches
    tb(0.42, 3.15, 1.5, 0.35, "Red Laser (650 nm)", 9, True)
    tb(0.42, 2.75, 1.55, 0.85, "Red channel:\nhorizontal end-face\nreference input", 7.5)
    tb(4.95, 6.55, 1.4, 0.45, "Step-index POF\n(9 cm rigid rod)", 8, align=PP_ALIGN.CENTER)
    tb(0.95, 2.45, 1.25, 0.65, "Left End-Face\n(Axial Reference Input)", 7.5)
    tb(7.95, 2.45, 1.25, 0.65, "Right End-Face\n(Output)", 7.5, align=PP_ALIGN.RIGHT)
    tb(5.55, 6.2, 1.45, 0.55, "Side-polished Region\n(Side Illumination / Challenge Input)", 7.5, align=PP_ALIGN.CENTER)
    tb(1.45, 7.35, 1.45, 0.35, "Green Laser (520 nm)", 9, True)
    tb(1.45, 6.95, 1.55, 0.75, "Green channel:\noblique side-illumination\nchallenge input", 7.5)
    tb(2.85, 6.05, 0.95, 0.35, "Beam Expander", 7.5, align=PP_ALIGN.CENTER)
    tb(3.65, 5.55, 0.95, 0.45, "Collimating Lens", 7.5, align=PP_ALIGN.CENTER)
    tb(4.55, 4.75, 1.05, 0.55, "SLM\n(Spatial Light Modulator)", 7.5, align=PP_ALIGN.CENTER)
    tb(6.05, 4.05, 0.95, 0.45, "Focusing Lens", 7.5, align=PP_ALIGN.CENTER)
    tb(9.6, 5.5, 1.0, 0.3, "CMOS Camera", 8, True, PP_ALIGN.CENTER)
    tb(9.85, 3.85, 0.65, 0.45, "Speckle\nOutput", 7.5, align=PP_ALIGN.CENTER)

    # Panel (b) title + labels overlaid near bottom picture
    tb(0.55, 1.75, 3.4, 0.35, "Magnified View: Side-polished Region", 9, True)
    tb(0.65, 0.58, 0.75, 0.5, "Cladding\n(Lower n)", 7, align=PP_ALIGN.CENTER)
    tb(2.95, 1.05, 0.6, 0.4, "Core\n(Higher n)", 7, align=PP_ALIGN.CENTER)
    tb(1.6, 1.78, 1.25, 0.4, "Polishing window\n(enhanced side coupling)", 7, align=PP_ALIGN.CENTER)
    tb(2.45, 1.12, 1.1, 0.4, "Side-illuminated\nscattering volume", 7)

    # Panel letters
    tb(0.42, 4.55, 0.35, 0.3, "(a)", 12, True)
    tb(0.42, 1.58, 0.35, 0.28, "(b)", 12, True)

    # Panel (c) — generate small PNG from matplotlib for placement
    fig, ax = plt.subplots(figsize=(3.2, 2.6), facecolor="white")
    draw_panel_c(ax, np.random.default_rng(7))
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=180, facecolor="white", bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, emu(9.55), emu(4.55), width=emu(3.35))
    tb(9.55, 7.12, 0.35, 0.28, "(c)", 12, True)

    # Scale bar (editable lines as thin rects)
    ysc = 2.2
    x0, x3 = 1.0, 7.6
    seg = [(x0, x0 + 5 / 9 * (x3 - x0)), (x0 + 5 / 9 * (x3 - x0), x0 + 6 / 9 * (x3 - x0)), (x0 + 6 / 9 * (x3 - x0), x3)]
    for (a, b) in seg:
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(a), emu(ysc), emu(b - a), emu(0.035))
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(26, 39, 68)
        r.line.fill.background()
    tb(3.25, 1.85, 0.55, 0.22, "5 cm", 7.5, align=PP_ALIGN.CENTER)
    tb(5.15, 1.85, 0.55, 0.22, "1 cm", 7.5, align=PP_ALIGN.CENTER)
    tb(6.85, 1.85, 0.55, 0.22, "3 cm", 7.5, align=PP_ALIGN.CENTER)
    tb(4.35, 1.58, 1.5, 0.22, "Total: 9 cm", 8.5, True, PP_ALIGN.CENTER)

    # Legend box
    leg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(9.5), emu(0.45), emu(3.35), emu(1.25))
    leg.fill.solid()
    leg.fill.fore_color.rgb = RGBColor(245, 245, 245)
    leg.line.color.rgb = RGBColor(176, 190, 197)
    tb(9.75, 1.42, 2.8, 0.22, "Red channel — (Reference, end-face axial)", 7.5, True)
    tb(9.75, 1.15, 2.8, 0.22, "Green channel — (Signal, side illumination)", 7.5, True)
    tb(9.75, 0.88, 2.8, 0.2, "POF Core (Higher n)", 7.5)
    tb(9.75, 0.65, 2.8, 0.2, "POF Cladding (Lower n)", 7.5)

    outp = OUT / f"{BASE}_semieditable.pptx"
    prs.save(str(outp))
    print("Wrote", outp)


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    compose_figure()
    _compose_annotations_svg_only()
    build_pptx_semieditable()


if __name__ == "__main__":
    main()
