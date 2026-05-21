#!/usr/bin/env python3
"""
Build Fig. 2 — dual-channel optical setup (9 cm POF: 5 + 1 + 3 cm).

Outputs (canonical):
  figures/paper/Fig2_setup/Fig2_optical_setup_9cm_editable.pptx  — native editable shapes + text
  figures/paper/Fig2_setup/Fig2_optical_setup_9cm.{svg,pdf,png} — matplotlib vector/raster

Run from repo root:
  pip install python-pptx
  python3 scripts/paper_figures/build_fig2_optical_setup.py
"""
from __future__ import annotations

import math
import sys
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib.colorbar import ColorbarBase
from matplotlib.patches import ConnectionPatch, FancyBboxPatch, Rectangle, Circle, Ellipse
from matplotlib.transforms import Affine2D
import numpy as np

SCRIPT_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRIPT_DIR.parents[1]
sys.path.insert(0, str(SCRIPT_DIR.parent))

OUT_DIR = REPO_ROOT / "figures" / "paper" / "Fig2_setup"
BASE = "Fig2_optical_setup_9cm"

# --- Colors (match brief) ---
C_BG = "#FFFFFF"
C_LABEL = "#1A2744"
C_RED = "#C62828"
C_RED_SOFT = "#E57373"
C_GREEN = "#2E7D32"
C_GREEN_SOFT = "#81C784"
C_GREEN_BEAM = "#4CAF50"
C_CORE = "#90CAF9"
C_CLAD = "#CFD8DC"
C_SLM_BG = "#263238"
C_SLM_SCR = "#1B5E20"
C_PURPLE_DASH = "#7B1FA2"
C_GREEN_DASH = "#388E3C"
C_LEGEND_BOX = "#F5F5F5"
DPI_PNG = 600


def _laser_body_mpl(ax, xy, w, h, angle_deg=0, z=10):
    """Layered pseudo-3D laser module."""
    rad = math.radians(angle_deg)
    c, s = math.cos(rad), math.sin(rad)
    R = np.array([[c, -s], [s, c]])
    corners = np.array([[-w / 2, -h / 2], [w / 2, -h / 2], [w / 2, h / 2], [-w / 2, h / 2]])
    rot = (R @ corners.T).T + np.array(xy)
    poly = plt.Polygon(rot, closed=True, facecolor="#2B2B2B", edgecolor="#0D0D0D", linewidth=0.8, zorder=z)
    ax.add_patch(poly)
    # top highlight
    hi = rot.copy()
    hi[:, 1] *= 0.92
    hi[2:, 1] += h * 0.08
    ax.add_patch(plt.Polygon(hi, closed=True, facecolor="#4A4A4A", edgecolor="none", zorder=z + 0.1, alpha=0.9))
    # fins
    for off in (-w * 0.35, w * 0.35):
        fx = xy[0] + off * c
        fy = xy[1] + off * s
        fins = FancyBboxPatch(
            (fx - w * 0.06, fy - h * 0.55), w * 0.12, h * 1.1,
            boxstyle="round,pad=0.01,rounding_size=0.02",
            facecolor="#1A1A1A", edgecolor="#333333", linewidth=0.4, zorder=z - 1,
        )
        ax.add_patch(fins)


def _lens_mpl(ax, cx, cy, rw, rh, ang_deg, z=8):
    ring = Ellipse((cx, cy), rw * 2, rh * 2, angle=ang_deg, facecolor="#ECEFF1", edgecolor="#546E7A", linewidth=1.0, zorder=z)
    ax.add_patch(ring)
    ax.add_patch(Ellipse((cx, cy), rw * 1.35, rh * 1.35, angle=ang_deg, facecolor="#E3F2FD", edgecolor="#78909C", linewidth=0.6, zorder=z + 0.1))


def _slm_mpl(ax, cx, cy, w, h, ang_deg, z=9):
    trans = Affine2D().rotate_deg_around(cx, cy, ang_deg) + ax.transData
    ax.add_patch(Rectangle(
        (cx - w / 2, cy - h / 2), w, h, transform=trans,
        facecolor=C_SLM_BG, edgecolor="#000000", linewidth=1.2, zorder=z,
    ))
    ax.add_patch(Rectangle(
        (cx - w * 0.36, cy - h * 0.31), w * 0.72, h * 0.62, transform=trans,
        facecolor=C_SLM_SCR, edgecolor="#004D40", linewidth=0.6, zorder=z + 1,
    ))


def _camera_mpl(ax, cx, cy, scale=1.0, z=10):
    w, h = 5.5 * scale, 4.2 * scale
    body = FancyBboxPatch(
        (cx - w / 2, cy - h / 2), w, h,
        boxstyle="round,pad=0.02,rounding_size=0.35",
        facecolor="#263238", edgecolor="#000000", linewidth=0.9, zorder=z,
    )
    ax.add_patch(body)
    ax.add_patch(Circle((cx + w * 0.25, cy), 1.35 * scale, facecolor="#37474F", edgecolor="#000000", linewidth=0.7, zorder=z + 1))
    ax.add_patch(Circle((cx + w * 0.25, cy), 0.65 * scale, facecolor="#0D47A1", edgecolor="#000000", linewidth=0.4, zorder=z + 2))


def build_matplotlib_figure():
    """Publication-grade SVG/PDF/PNG."""
    fig = plt.figure(figsize=(11.5, 8.2), facecolor=C_BG, dpi=DPI_PNG)
    fig.patch.set_facecolor(C_BG)

    # --- Panel (a) main axes (model coords) ---
    ax = fig.add_axes([0.085, 0.30, 0.58, 0.58])
    ax.set_facecolor(C_BG)
    ax.set_xlim(0, 100)
    ax.set_ylim(10, 92)
    ax.set_aspect("equal")
    ax.axis("off")

    # Fiber geometry: x_left to x_right, total 9 cm => 5+1+3
    y_clad_b, y_clad_t = 36.0, 48.5
    y_core_b, y_core_t = 38.2, 46.2
    x_left, x_right = 17.5, 87.0
    L = x_right - x_left
    x_pol0 = x_left + L * (5.0 / 9.0)
    x_pol1 = x_left + L * (6.0 / 9.0)

    # Cladding + core
    clad = FancyBboxPatch(
        (x_left, y_clad_b), x_right - x_left, y_clad_t - y_clad_b,
        boxstyle="round,pad=0,rounding_size=1.2",
        facecolor=C_CLAD, edgecolor="#90A4AE", linewidth=1.0, zorder=2, alpha=0.95,
    )
    ax.add_patch(clad)
    core = FancyBboxPatch(
        (x_left, y_core_b), x_right - x_left, y_core_t - y_core_b,
        boxstyle="round,pad=0,rounding_size=0.85",
        facecolor=C_CORE, edgecolor="#1565C0", linewidth=0.9, zorder=3, alpha=0.98,
    )
    ax.add_patch(core)

    # Side-polished window (1 cm)
    polish = Rectangle(
        (x_pol0, y_clad_b), x_pol1 - x_pol0, y_clad_t - y_clad_b,
        facecolor=C_GREEN_SOFT, linewidth=0, zorder=4, alpha=0.42,
    )
    ax.add_patch(polish)
    ax.add_patch(Rectangle((x_pol0, y_clad_b), x_pol1 - x_pol0, y_clad_t - y_clad_b, fill=False, edgecolor=C_GREEN, linewidth=1.8, zorder=5, linestyle="-"))

    # Green glow at coupling
    glow = Ellipse(((x_pol0 + x_pol1) / 2, y_core_t + 0.8), (x_pol1 - x_pol0) * 1.15, 5.5, facecolor="#A5D6A7", alpha=0.55, zorder=6, edgecolor="none")
    ax.add_patch(glow)

    # Red beam (horizontal)
    ax.plot([7.5, x_left], [42.0, 42.0], color=C_RED, solid_capstyle="round", linewidth=3.8, alpha=0.85, zorder=7)
    ax.plot([7.5, x_left], [42.0, 42.0], color=C_RED_SOFT, solid_capstyle="round", linewidth=1.4, alpha=0.55, zorder=8)

    _laser_body_mpl(ax, (5.0, 42.0), 7.0, 5.2, angle_deg=0)

    # Green path: collinear segments, oblique
    s_start = np.array([23.0, 76.0])
    s_hit = np.array([(x_pol0 + x_pol1) / 2, y_clad_t + 0.4])
    direction = s_hit - s_start
    direction = direction / np.linalg.norm(direction)
    total_len = np.linalg.norm(s_hit - s_start) - 1.2
    ts = np.linspace(2.0, total_len, 80)
    beam_pts = np.array([s_start + t * direction for t in ts])
    ax.plot(beam_pts[:, 0], beam_pts[:, 1], color=C_GREEN_BEAM, linewidth=4.0, alpha=0.38, zorder=7, solid_capstyle="round")
    ax.plot(beam_pts[:, 0], beam_pts[:, 1], color="#C8E6C9", linewidth=1.5, alpha=0.65, zorder=8)

    # Optics along path
    def pt_at(frac: float):
        return s_start + (2.0 + frac * (total_len - 2.0)) * direction

    ang = math.degrees(math.atan2(direction[1], direction[0]))
    _laser_body_mpl(ax, tuple(s_start - direction * 3.2), 6.2, 4.8, angle_deg=ang)

    p_be = pt_at(0.08)
    p_cl = pt_at(0.22)
    p_slm = pt_at(0.42)
    p_fo = pt_at(0.62)
    _lens_mpl(ax, p_be[0], p_be[1], 2.1, 0.95, ang_deg=ang)
    _lens_mpl(ax, p_cl[0], p_cl[1], 2.0, 0.9, ang_deg=ang)
    _slm_mpl(ax, p_slm[0], p_slm[1], 6.0, 4.0, ang_deg=ang - 8)
    _lens_mpl(ax, p_fo[0], p_fo[1], 2.3, 1.0, ang_deg=ang)

    # CMOS + speckle
    _camera_mpl(ax, 93.5, 42.5, scale=0.95)
    rng = np.random.default_rng(42)
    sp = rng.normal(0, 1, (32, 32))
    ax_in = ax.inset_axes([0.785, 0.38, 0.11, 0.14], transform=ax.transAxes, zorder=20)
    ax_in.imshow(sp, cmap="gray", vmin=-2, vmax=2, aspect="equal")
    ax_in.axis("off")

    # --- Labels panel (a) ---
    def txt(x, y, s, fs=7.5, ha="left", va="bottom", bold=False):
        ax.text(x, y, s, fontsize=fs, color=C_LABEL, ha=ha, va=va, fontweight="bold" if bold else "normal",
                fontfamily="sans-serif")

    txt(1.0, 33.5, "Red Laser (650 nm)", fs=7.8, bold=True)
    txt(1.0, 31.0, "Red channel:\nhorizontal end-face\nreference input", fs=7.2)
    txt(48, 49.5, "Step-index POF\n(9 cm rigid rod)", fs=7.5, ha="center")
    txt(x_left - 0.5, y_clad_b - 3.8, "Left End-Face\n(Axial Reference Input)", fs=7.2, ha="left")
    txt(x_right - 2, y_clad_b - 3.8, "Right End-Face\n(Output)", fs=7.2, ha="right")
    txt((x_pol0 + x_pol1) / 2, y_clad_t + 6.5, "Side-polished Region\n(Side Illumination / Challenge Input)", fs=7.2, ha="center")
    txt(15.0, 79.0, "Green Laser (520 nm)", fs=7.8, bold=True)
    txt(15.0, 76.0, "Green channel:\noblique side-illumination\nchallenge input", fs=7.2)
    txt(p_be[0] + 2, p_be[1] + 4, "Beam Expander", fs=7.0, ha="center")
    txt(p_cl[0] + 2.2, p_cl[1] + 3.5, "Collimating\nLens", fs=7.0, ha="center")
    txt(p_slm[0], p_slm[1] - 7.5, "SLM\n(Spatial Light Modulator)", fs=7.0, ha="center")
    txt(p_fo[0] + 2.5, p_fo[1] - 1.0, "Focusing\nLens", fs=7.0, ha="center")
    txt(88, 49, "CMOS Camera", fs=7.5, bold=True)
    txt(91, 33, "Speckle\nOutput", fs=7.0, ha="center")

    # Scale: 5 cm | 1 cm | 3 cm under fiber
    y_br = 25.0
    h_br = 1.0
    ax.plot([x_left, x_pol0], [y_br, y_br], color=C_LABEL, linewidth=1.0)
    ax.plot([x_pol0, x_pol1], [y_br, y_br], color=C_LABEL, linewidth=1.0)
    ax.plot([x_pol1, x_right], [y_br, y_br], color=C_LABEL, linewidth=1.0)
    for x0, x1, lab in ((x_left, x_pol0, "5 cm"), (x_pol0, x_pol1, "1 cm"), (x_pol1, x_right, "3 cm")):
        mx = (x0 + x1) / 2
        ax.plot([x0, x0], [y_br - h_br, y_br + h_br], color=C_LABEL, linewidth=0.9)
        ax.plot([x1, x1], [y_br - h_br, y_br + h_br], color=C_LABEL, linewidth=0.9)
        ax.text(mx, y_br - 3.2, lab, ha="center", va="top", fontsize=7.0, color=C_LABEL)
    ax.text((x_left + x_right) / 2, y_br - 7.5, "Total: 9 cm", ha="center", va="top", fontsize=7.2, color=C_LABEL, fontweight="bold")

    ax.text(0.02, 0.98, "(a)", transform=ax.transAxes, fontsize=11, fontweight="bold", va="top", color=C_LABEL)

    # --- Panel (b): magnified inset ---
    ax_b = fig.add_axes([0.055, 0.065, 0.30, 0.205])
    ax_b.set_xlim(0, 100)
    ax_b.set_ylim(0, 40)
    ax_b.axis("off")
    ax_b.patch.set_facecolor(C_BG)
    bbox = FancyBboxPatch(
        (2, 2), 96, 36, boxstyle="round,pad=0,rounding_size=2.5",
        facecolor="none", edgecolor=C_PURPLE_DASH, linewidth=2.0, linestyle=(0, (5, 3)),
    )
    ax_b.add_patch(bbox)
    ax_b.set_title("Magnified View: Side-polished Region", fontsize=8.5, fontweight="bold", color=C_LABEL, pad=6)

    yb0, yb1 = 14.0, 26.0
    xc0, xc1 = 12.0, 88.0
    ax_b.add_patch(FancyBboxPatch((xc0, yb0), xc1 - xc0, yb1 - yb0, boxstyle="round,pad=0.4,rounding_size=2.8",
                                  facecolor=C_CLAD, edgecolor="#90A4AE", linewidth=1.0, alpha=0.95))
    ax_b.add_patch(FancyBboxPatch((xc0 + 1.5, yb0 + 2.2), xc1 - xc0 - 3, yb1 - yb0 - 4.4, boxstyle="round,pad=0.3,rounding_size=1.8",
                                  facecolor=C_CORE, edgecolor="#1565C0", linewidth=0.8))
    pw0, pw1 = 44.0, 58.0
    ax_b.add_patch(Rectangle((pw0, yb1 - 1.5), pw1 - pw0, 8.0, facecolor="#A5D6A7", alpha=0.65, edgecolor=C_GREEN, linewidth=1.0))
    ax_b.plot([51, 51], [yb1 + 7, 34], color=C_GREEN_BEAM, linewidth=3.5, alpha=0.5)
    ax_b.add_patch(Ellipse((51, 22), 14, 8, facecolor="#C8E6C9", alpha=0.6, edgecolor="none"))
    ax_b.text(50, 10, "Cladding\n(Lower n)", ha="center", fontsize=7, color=C_LABEL)
    ax_b.text(75, 19, "Core\n(Higher n)", ha="center", fontsize=7, color=C_LABEL)
    ax_b.text(51, 32.5, "Polishing window\n(enhanced side coupling)", ha="center", fontsize=7, color=C_LABEL)
    ax_b.text(62, 24, "Side-illuminated\nscattering volume", ha="left", fontsize=7, color=C_LABEL)
    ax_b.text(0.02, 0.98, "(b)", transform=ax_b.transAxes, fontsize=11, fontweight="bold", va="top", color=C_LABEL)

    # --- Panel (c): SLM examples ---
    ax_c = fig.add_axes([0.708, 0.52, 0.265, 0.315])
    ax_c.set_xlim(0, 100)
    ax_c.set_ylim(0, 100)
    ax_c.axis("off")
    cb2 = FancyBboxPatch(
        (3, 3), 94, 94, boxstyle="round,pad=0,rounding_size=2.5",
        facecolor="none", edgecolor=C_GREEN_DASH, linewidth=2.0, linestyle=(0, (5, 3)),
    )
    ax_c.add_patch(cb2)
    ax_c.set_title("SLM Letter Image Input (Examples)", fontsize=8.5, fontweight="bold", color=C_LABEL, pad=8)

    letters = [("A", 12), ("B", 32), ("C", 52), ("Z", 72)]
    rng2 = np.random.default_rng(1)
    for ch, x0 in letters:
        ax_c.add_patch(Rectangle((x0, 58), 16, 18, facecolor="#000000", edgecolor="#333333", linewidth=0.8))
        ax_c.text(x0 + 8, 67, ch, ha="center", va="center", fontsize=14, color="white", fontweight="bold", fontfamily="sans-serif")
        yy, xx = np.mgrid[0:16, 0:16]
        ph = np.sin(xx * 0.45 + yy * 0.38 + rng2.normal(0, 0.2, (16, 16))) * 0.5 + rng2.normal(0, 0.05, (16, 16))
        ax_c.imshow(ph, extent=(x0, x0 + 16, 40, 56), cmap="viridis", origin="lower", vmin=-1, vmax=1, clip_on=True)

    ax_c.text(50, 8, "26 challenge patterns (A–Z)", ha="center", fontsize=7.5, color=C_LABEL)
    ax_c.text(0.02, 0.98, "(c)", transform=ax_c.transAxes, fontsize=11, fontweight="bold", va="top", color=C_LABEL)

    # Phase colorbar
    cax = fig.add_axes([0.965, 0.58, 0.014, 0.22])
    ColorbarBase(cax, cmap=plt.cm.viridis, orientation="vertical")
    cax.set_title("Phase", fontsize=7, color=C_LABEL, pad=4)
    cax.set_ylabel("2π", fontsize=6.5, color=C_LABEL, labelpad=2)
    cax.yaxis.set_ticks_position("right")

    # Connector (c) -> SLM in figure coordinates
    slm_disp = ax.transData.transform(p_slm)
    inv = fig.transFigure.inverted()
    slm_fig = inv.transform(slm_disp)
    c_box = [0.708, 0.52, 0.265, 0.315]
    start = (c_box[0] + c_box[2] * 0.15, c_box[1])
    con = ConnectionPatch(
        xyA=start, coordsA=fig.transFigure, axesA=None,
        xyB=(slm_fig[0] + 0.02, slm_fig[1] - 0.02), coordsB=fig.transFigure, axesB=None,
        linestyle=(0, (4, 3)), linewidth=1.2, color=C_GREEN_BEAM, zorder=100,
    )
    fig.add_artist(con)

    # --- Legend ---
    ax_l = fig.add_axes([0.68, 0.07, 0.275, 0.16])
    ax_l.set_xlim(0, 1)
    ax_l.set_ylim(0, 1)
    ax_l.axis("off")
    box = FancyBboxPatch(
        (0.03, 0.08), 0.94, 0.84, boxstyle="round,pad=0.02,rounding_size=0.03",
        facecolor=C_LEGEND_BOX, edgecolor="#B0BEC5", linewidth=1.0,
    )
    ax_l.add_patch(box)
    ax_l.plot([0.1, 0.2], [0.78, 0.78], color=C_RED, linewidth=3, solid_capstyle="round")
    ax_l.text(0.24, 0.78, "Red channel", fontsize=7.5, color=C_LABEL, va="center", fontweight="bold")
    ax_l.text(0.24, 0.65, "(Reference, end-face axial)", fontsize=6.8, color=C_LABEL, va="center")
    ax_l.plot([0.1, 0.2], [0.48, 0.48], color=C_GREEN_BEAM, linewidth=3, solid_capstyle="round", alpha=0.85)
    ax_l.text(0.24, 0.48, "Green channel", fontsize=7.5, color=C_LABEL, va="center", fontweight="bold")
    ax_l.text(0.24, 0.35, "(Signal, side illumination)", fontsize=6.8, color=C_LABEL, va="center")
    ax_l.add_patch(Rectangle((0.08, 0.14), 0.12, 0.08, facecolor=C_CORE, edgecolor="#1565C0", linewidth=0.5))
    ax_l.text(0.24, 0.18, "POF Core (Higher n)", fontsize=7, color=C_LABEL, va="center")
    ax_l.add_patch(Rectangle((0.08, 0.02), 0.12, 0.07, facecolor=C_CLAD, edgecolor="#90A4AE", linewidth=0.5))
    ax_l.text(0.24, 0.055, "POF Cladding (Lower n)", fontsize=7, color=C_LABEL, va="center")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for ext in ("svg", "pdf", "png"):
        p = OUT_DIR / f"{BASE}.{ext}"
        kwargs = {"dpi": DPI_PNG} if ext == "png" else {}
        fig.savefig(p, format=ext, facecolor=C_BG, **kwargs)
    plt.close(fig)
    print("Wrote", OUT_DIR / f"{BASE}.svg", "etc.")


def build_pptx():
    """Editable PowerPoint (native shapes + text)."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Emu, Inches, Pt
    except ImportError as e:
        print("python-pptx required: pip install python-pptx", file=sys.stderr)
        raise

    # EMU helper
    def emu(x_inch: float) -> int:
        return int(round(x_inch * 914400))

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = emu(13.333)  # 16:9
    prs.slide_height = emu(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Slide origin and scale (match mpl layout loosely)
    W, H = 13.333, 7.5
    ox, oy = 0.55, 0.45

    def lbl(x, y, w, h, text, font_pt=7, bold=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(emu(ox + x), emu(oy + y), emu(w), emu(h))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.name = "Arial"
        p.font.size = Pt(font_pt)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(26, 39, 68)
        return tb

    def rect(x, y, w, h, fill_rgb, line_rgb=None, line_w=Pt(0.5)):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(ox + x), emu(oy + y), emu(w), emu(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(*fill_rgb)
        if line_rgb:
            sh.line.color.rgb = RGBColor(*line_rgb)
            sh.line.width = Pt(1)
        else:
            sh.line.fill.background()
        return sh

    def round_rect(x, y, w, h, fill_rgb, outline_rgb, dash=False):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(ox + x), emu(oy + y), emu(w), emu(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(*fill_rgb)
        sh.line.color.rgb = RGBColor(*outline_rgb)
        sh.line.width = Pt(1.5)
        if dash:
            sh.line.dash_style = 2
        return sh

    def line_xy(x1, y1, x2, y2, rgb, w_pt=3.0):
        cx, cy = (x1 + x2) / 2, (y1 + y2) / 2
        dx, dy = x2 - x1, y2 - y1
        ln = (dx * dx + dy * dy) ** 0.5
        if ln < 1e-6:
            return None
        sh = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            emu(ox + cx - ln / 2), emu(oy + cy - 0.03), emu(ln), emu(0.06),
        )
        sh.rotation = math.degrees(math.atan2(dy, dx))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(*rgb)
        sh.line.fill.background()
        return sh

    # Background white
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

    # Fiber (cladding + core)
    y_f = 3.35
    rect(1.15, y_f - 0.12, 7.45, 0.50, (207, 215, 216), (144, 164, 174))
    rect(1.15, y_f - 0.02, 7.45, 0.30, (144, 202, 249), (21, 101, 192))

    x_pol0 = 1.15 + 7.45 * (5 / 9)
    x_pol1 = 1.15 + 7.45 * (6 / 9)
    sh_pol = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(ox + x_pol0), emu(oy + y_f - 0.12), emu(x_pol1 - x_pol0), emu(0.50))
    sh_pol.fill.solid()
    sh_pol.fill.fore_color.rgb = RGBColor(129, 199, 132)
    sh_pol.fill.transparency = 0.45
    sh_pol.line.color.rgb = RGBColor(46, 125, 50)
    sh_pol.line.width = Pt(2)

    # Red beam
    line_xy(0.35, y_f + 0.13, 1.12, y_f + 0.13, (198, 40, 40), 5)
    # Red laser body
    rect(0.12, y_f - 0.08, 0.38, 0.42, (43, 43, 43), (26, 26, 26))

    # Green path (single oblique bar as composite segments) — approximate with line
    gx0, gy0 = 2.36, 5.85
    gx1, gy1 = (x_pol0 + x_pol1) / 2, y_f + 0.52
    nseg = 12
    for i in range(nseg):
        t0, t1 = i / nseg, (i + 1) / nseg
        px0, py0 = gx0 + t0 * (gx1 - gx0), gy0 + t0 * (gy1 - gy0)
        px1, py1 = gx0 + t1 * (gx1 - gx0), gy0 + t1 * (gy1 - gy0)
        line_xy(px0, py0, px1, py1, (76, 175, 80), 5)

    rect(1.95, 5.45, 0.42, 0.38, (43, 43, 43), (26, 26, 26))

    # Lenses (ovals)
    for (cx, cy) in ((3.05, 5.15), (3.65, 4.68), (5.35, 3.82)):
        el = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu(ox + cx - 0.22), emu(oy + cy - 0.12), emu(0.44), emu(0.24))
        el.fill.solid()
        el.fill.fore_color.rgb = RGBColor(227, 242, 253)
        el.line.color.rgb = RGBColor(84, 110, 122)

    # SLM
    slm = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(ox + 4.08), emu(oy + 4.05), emu(0.78), emu(0.55))
    slm.fill.solid()
    slm.fill.fore_color.rgb = RGBColor(38, 50, 56)
    slm.line.width = Pt(1.2)
    inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(ox + 4.18), emu(oy + 4.15), emu(0.55), emu(0.38))
    inner.fill.solid()
    inner.fill.fore_color.rgb = RGBColor(27, 94, 32)

    # Camera + speckle
    round_rect(9.05, y_f - 0.12, 0.55, 0.48, (38, 50, 56), (0, 0, 0))
    rect(9.78, y_f + 0.02, 0.38, 0.34, (224, 224, 224), (80, 80, 80))

    # Labels — panel (a)
    lbl(0.02, 2.25, 1.6, 0.5, "Red Laser (650 nm)", 8, bold=True)
    lbl(0.02, 1.85, 1.85, 0.85, "Red channel:\nhorizontal end-face\nreference input", 7)
    lbl(4.6, 4.15, 1.4, 0.45, "Step-index POF\n(9 cm rigid rod)", 7.5, align=PP_ALIGN.CENTER)
    lbl(1.0, 2.55, 1.4, 0.65, "Left End-Face\n(Axial Reference Input)", 7)
    lbl(7.95, 2.55, 1.3, 0.65, "Right End-Face\n(Output)", 7, align=PP_ALIGN.RIGHT)
    pcx = (x_pol0 + x_pol1) / 2
    lbl(pcx - 0.45, 4.05, 1.3, 0.55, "Side-polished Region\n(Side Illumination / Challenge Input)", 7, align=PP_ALIGN.CENTER)
    lbl(1.62, 6.35, 1.4, 0.35, "Green Laser (520 nm)", 8, bold=True)
    lbl(1.62, 5.9, 1.55, 0.75, "Green channel:\noblique side-illumination\nchallenge input", 7)
    lbl(2.78, 5.6, 0.85, 0.35, "Beam Expander", 7, align=PP_ALIGN.CENTER)
    lbl(3.42, 5.05, 0.75, 0.45, "Collimating\nLens", 7, align=PP_ALIGN.CENTER)
    lbl(4.0, 4.2, 0.95, 0.55, "SLM\n(Spatial Light Modulator)", 7, align=PP_ALIGN.CENTER)
    lbl(5.05, 3.55, 0.85, 0.45, "Focusing\nLens", 7, align=PP_ALIGN.CENTER)
    lbl(8.82, 4.05, 0.9, 0.3, "CMOS Camera", 7.5, bold=True)
    lbl(9.15, 2.7, 0.55, 0.45, "Speckle\nOutput", 7, align=PP_ALIGN.CENTER)

    # Scale bar
    ysc = 2.95
    line_xy(1.15, ysc, x_pol0, ysc, (26, 39, 68), 1.5)
    line_xy(x_pol0, ysc, x_pol1, ysc, (26, 39, 68), 1.5)
    line_xy(x_pol1, ysc, 8.6, ysc, (26, 39, 68), 1.5)
    lbl(3.2, 2.15, 0.5, 0.25, "5 cm", 7, align=PP_ALIGN.CENTER)
    lbl(5.15, 2.15, 0.5, 0.25, "1 cm", 7, align=PP_ALIGN.CENTER)
    lbl(7.15, 2.15, 0.5, 0.25, "3 cm", 7, align=PP_ALIGN.CENTER)
    lbl(4.6, 1.78, 2.0, 0.22, "Total: 9 cm", 7.5, bold=True, align=PP_ALIGN.CENTER)

    lbl(0.15, 6.9, 0.35, 0.3, "(a)", 11, bold=True)

    # Panel (b)
    bx, by, bw, bh = 0.12, 0.38, 3.45, 1.4
    round_rect(bx, by, bw, bh, (255, 255, 255), (123, 31, 162), dash=True)
    lbl(bx + 0.35, by + 1.15, 2.8, 0.35, "Magnified View: Side-polished Region", 8.5, bold=True)
    rect(bx + 0.35, by + 0.42, 2.85, 0.52, (207, 215, 216), (144, 164, 174))
    rect(bx + 0.5, by + 0.52, 2.55, 0.32, (144, 202, 249), (21, 101, 192))
    rect(bx + 1.4, by + 0.88, 0.75, 0.22, (165, 214, 167), (46, 125, 50))
    line_xy(bx + 1.77, by + 1.35, bx + 1.77, by + 1.02, (76, 175, 80), 4)
    lbl(bx + 1.0, by + 0.12, 0.8, 0.5, "Cladding\n(Lower n)", 7, align=PP_ALIGN.CENTER)
    lbl(bx + 2.25, by + 0.55, 0.55, 0.45, "Core\n(Higher n)", 7, align=PP_ALIGN.CENTER)
    lbl(bx + 1.25, by + 1.18, 1.2, 0.4, "Polishing window\n(enhanced side coupling)", 7, align=PP_ALIGN.CENTER)
    lbl(bx + 0.12, by + 1.05, 0.32, 0.28, "(b)", 11, bold=True)

    # Panel (c)
    cx, cy, cw, ch = 9.65, 4.6, 3.2, 2.45
    round_rect(cx, cy, cw, ch, (255, 255, 255), (56, 142, 60), dash=True)
    lbl(cx + 0.35, cy + 2.08, 2.6, 0.35, "SLM Letter Image Input (Examples)", 8.5, bold=True)
    for i, letter in enumerate(["A", "B", "C", "Z"]):
        lx = cx + 0.25 + i * 0.72
        rect(lx, cy + 1.28, 0.48, 0.55, (0, 0, 0), (50, 50, 50))
        t = slide.shapes.add_textbox(emu(ox + lx), emu(oy + cy + 1.38), emu(0.48), emu(0.4))
        t.text_frame.paragraphs[0].text = letter
        t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        pf = t.text_frame.paragraphs[0].font
        pf.size = Pt(14)
        pf.bold = True
        pf.color.rgb = RGBColor(255, 255, 255)
        pf.name = "Arial"
        rect(lx, cy + 0.62, 0.48, 0.58, (69, 90, 100), (40, 53, 60))
    lbl(cx + 0.35, cy + 0.12, 2.6, 0.28, "26 challenge patterns (A–Z)", 7.5, align=PP_ALIGN.CENTER)
    lbl(cx + 0.05, cy + 2.35, 0.32, 0.28, "(c)", 11, bold=True)

    # Dashed connector (c) to SLM — polyline via small segments
    line_xy(cx + 0.4, cy + 2.45, 4.9, 4.85, (76, 175, 80), 2)

    # Legend
    lx, ly = 9.6, 0.45
    round_rect(lx, ly, 3.25, 1.22, (245, 245, 245), (176, 190, 197))
    line_xy(lx + 0.18, ly + 0.92, lx + 0.45, ly + 0.92, (198, 40, 40), 4)
    lbl(lx + 0.5, ly + 0.84, 2.4, 0.2, "Red channel", 7.5, bold=True)
    lbl(lx + 0.5, ly + 0.68, 2.4, 0.2, "(Reference, end-face axial)", 6.8)
    line_xy(lx + 0.18, ly + 0.54, lx + 0.45, ly + 0.54, (76, 175, 80), 4)
    lbl(lx + 0.5, ly + 0.46, 2.4, 0.2, "Green channel", 7.5, bold=True)
    lbl(lx + 0.5, ly + 0.30, 2.4, 0.2, "(Signal, side illumination)", 6.8)
    rect(lx + 0.12, ly + 0.16, 0.22, 0.09, (144, 202, 249), (21, 101, 192))
    lbl(lx + 0.38, ly + 0.12, 2.0, 0.14, "POF Core (Higher n)", 7)
    rect(lx + 0.12, ly + 0.02, 0.22, 0.08, (207, 215, 216), (144, 164, 174))
    lbl(lx + 0.38, ly + 0.0, 2.0, 0.12, "POF Cladding (Lower n)", 7)

    out_pptx = OUT_DIR / f"{BASE}_editable.pptx"
    prs.save(str(out_pptx))
    print("Wrote", out_pptx)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_matplotlib_figure()
    build_pptx()


if __name__ == "__main__":
    main()